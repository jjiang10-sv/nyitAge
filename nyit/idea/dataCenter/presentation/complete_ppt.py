from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import os
import subprocess

# ---------------------------
# HELPER FUNCTIONS from full_ppt.py
# ---------------------------

def add_bullet_slide(prs, title, bullets):
    """Add slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0

def add_table_slide(prs, title, data, col_widths=None):
    """Add slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank slide
    shapes = slide.shapes
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_shape.text = title
    rows, cols = len(data), len(data[0])
    table = shapes.add_table(rows, cols, Inches(0.5), Inches(1), Inches(9), Inches(2)).table
    for r in range(rows):
        for c in range(cols):
            table.cell(r, c).text = str(data[r][c])
            table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(12)
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

def latex_to_png(latex, filename, dpi=200):
    """Render LaTeX equation to PNG (requires pdflatex and ImageMagick convert)"""
    try:
        tex = r"""\documentclass[border=2pt]{standalone}
\usepackage{amsmath}
\begin{document}
%s
\end{document}""" % latex
        with open("temp.tex", "w") as f:
            f.write(tex)
        # Compile using pdflatex
        subprocess.run(["pdflatex", "-interaction=nonstopmode", "temp.tex"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        # Convert PDF to PNG
        subprocess.run(["convert", "-density", str(dpi), "temp.pdf", filename], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        # Cleanup
        for ext in ["aux","log","pdf","tex"]:
            if os.path.exists(f"temp.{ext}"):
                os.remove(f"temp.{ext}")
        return True
    except Exception as e:
        print(f"Warning: Could not generate LaTeX equation ({e}). Skipping equation slide.")
        return False

def add_image_slide(prs, title, image_path, width=Inches(8)):
    """Add slide with an image"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = title
    slide.shapes.add_picture(image_path, Inches(0.5), Inches(1), width=width)

def create_architecture_diagram():
    """Create architecture diagram using matplotlib"""
    fig, ax = plt.subplots(figsize=(8,4))
    ax.text(0.1, 0.8, "Traffic In", fontsize=12, fontweight="bold")
    ax.text(0.1, 0.5, "XDP/eBPF Pre-filter", fontsize=12, bbox=dict(facecolor='lightblue'))
    ax.arrow(1.5, 0.5, 1, 0, head_width=0.05, head_length=0.1)
    ax.text(2.7, 0.5, "Cilium Dataplane", fontsize=12, bbox=dict(facecolor='lightgreen'))
    ax.arrow(4.0, 0.5, 1, 0, head_width=0.05, head_length=0.1)
    ax.text(5.2, 0.5, "Pods + ML Autoscaler", fontsize=12, bbox=dict(facecolor='lightyellow'))
    ax.axis('off')
    plt.savefig("diagram.png", bbox_inches='tight')
    plt.close()

# ---------------------------
# CONTENT from ppt.py
# ---------------------------

slides_content = [
    {
        "title": "Cloud Center Security: Enhancing Kubernetes Security",
        "subtitle": "Cilium, eBPF, and DDoS Mitigation\nZhijun Jiang, Amin Milani Fard, NYIT, Vancouver, Canada"
    },
    {
        "title": "Motivation",
        "content": [
            "Kubernetes widely used for container orchestration.",
            "Cloud services are increasingly targeted by DDoS attacks.",
            "Current autoscaling is reactive and lacks traffic differentiation.",
            "Need security-aware, predictive autoscaling."
        ]
    },
    {
        "title": "Problem Statement",
        "content": [
            "Reactive scaling triggers on all traffic → may scale up for attacks.",
            "Network and security policies lag behind scaling events.",
            "Observability is fragmented → difficult to detect malicious patterns."
        ]
    },
    {
        "title": "Objectives",
        "content": [
            "Enhance Kubernetes autoscaling with security awareness.",
            "Integrate Cilium + eBPF for fast, programmable networking.",
            "Implement DDoS mitigation pipeline.",
            "Enable predictive scaling using ML traffic forecasts.",
            "Provide observability via Hubble + Grafana dashboards."
        ]
    },
    {
        "title": "Architecture Overview",
        "content": [
            "Edge Load Balancer / Cloud Scrubber",
            "XDP eBPF pre-filter → drop/rate-limit malicious traffic",
            "Cilium dataplane → enforces L3/L4/L7 policies",
            "Hubble → observability and flow logs",
            "Security-aware autoscaler → ML + Prometheus + policies"
        ]
    },
    {
        "title": "Components",
        "content": [
            "XDP / eBPF: Early packet filtering, drops malicious traffic",
            "Cilium: Pod-to-pod networking, policy enforcement",
            "Hubble: Observability, flow logs, metrics",
            "Predictive Autoscaler: Forecast traffic, scale pods securely",
            "DDoS Engine: Attack detection and mitigation actions"
        ]
    },
    {
        "title": "DDoS Mitigation Workflow",
        "content": [
            "1. Traffic enters cluster",
            "2. XDP pre-filter drops attack traffic",
            "3. Cilium enforces network policies",
            "4. Hubble observes flows → alerts DDoS engine",
            "5. Security-aware autoscaler scales only legitimate traffic"
        ]
    },
    {
        "title": "Predictive Scaling Model",
        "content": [
            "LSTM-based time series prediction",
            "Input: last 60 time steps of traffic metrics",
            "Output: next 12-step traffic forecast",
            "Integrates with Prometheus metrics and autoscaler",
            "Equation: ŷ_{t+1:t+12} = f(x_{t-59:t})"
        ]
    },
    {
        "title": "Key Features",
        "content": [
            "Real-time security-aware scaling",
            "Fast eBPF-based load balancing via Cilium",
            "Observability: per-pod traffic flow, status code tracking",
            "Predictive: anticipates spikes to avoid overloading",
            "Integrated DDoS mitigation"
        ]
    },
    {
        "title": "Experimental Setup",
        "content": [
            "Cluster: AWS EKS, 4 nodes (m5.xlarge)",
            "Workload: Go microservice, 5-min startup",
            "Traffic: Vegeta simulated (legitimate + attack mix)",
            "Monitoring: Prometheus + Grafana + Hubble",
            "Autoscalers: KEDA + PredictKube + custom ML controller"
        ]
    },
    {
        "title": "Results: Connection Stability",
        "content": [
            "Connection drops reduced by 72% with Cilium + autoscaler",
            "Predictive autoscaler reduced scaling lag by 54.8%",
            "ML forecast prevented unnecessary pod scaling during DDoS"
        ]
    },
    {
        "title": "Results: DDoS Mitigation",
        "content": [
            "XDP pre-filter dropped malicious requests immediately",
            "Cilium network policies restricted attack propagation",
            "Observability via Hubble → real-time flow detection"
        ]
    },
    {
        "title": "Related Work",
        "content": [
            "Kubernetes autoscaling: HPA, KEDA",
            "eBPF and Cilium for cloud security",
            "DDoS detection and mitigation in cloud",
            "Predictive scaling for container workloads"
        ]
    },
    {
        "title": "Lessons Learned",
        "content": [
            "Observability is critical → Grafana + Hubble",
            "Predictive scaling improves resource efficiency",
            "DDoS mitigation must integrate with scaling → avoid waste",
            "Cilium + eBPF ensures instant networking convergence"
        ]
    },
    {
        "title": "Future Work",
        "content": [
            "Hybrid predictive + reactive security-aware controller",
            "Multi-cloud deployment for resilience",
            "Chaos engineering to test autoscaler under attack",
            "Integrate anomaly detection for novel attacks"
        ]
    },
    {
        "title": "Conclusion",
        "content": [
            "Cloud center security requires integrated network + scaling + observability",
            "Cilium + eBPF + ML autoscaler improves performance and mitigates DDoS",
            "Security-aware scaling ensures efficient resource use",
            "Architecture is extendable and observable"
        ]
    },
    {
        "title": "References",
        "content": [
            "Kubernetes Event-Driven Autoscaling (KEDA): https://keda.sh/",
            "Cilium: eBPF Networking for Kubernetes: https://cilium.io/",
            "Hubble: Observability for Cilium: https://cilium.io/hubble/",
            "DDoS Attack Detection & Mitigation Survey, IEEE Access, 2022",
            "PredictKube: ML for Kubernetes Autoscaling: https://predictkube.com/",
            "eBPF Literature: BPF Performance and Applications, ACM 2020"
        ]
    }
]

# ---------------------------
# CREATE COMPLETE PRESENTATION
# ---------------------------

def main():
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = slides_content[0]["title"]
    slide.placeholders[1].text = slides_content[0]["subtitle"]
    
    # Add slides 1-4: Motivation through Objectives
    for s in slides_content[1:5]:
        add_bullet_slide(prs, s["title"], s["content"])
    
    # Add Architecture Overview (slide 5)
    add_bullet_slide(prs, slides_content[5]["title"], slides_content[5]["content"])
    
    # Add Architecture diagram
    create_architecture_diagram()
    add_image_slide(prs, "Cilium + KEDA Datapath Diagram", "diagram.png")
    
    # Add slides 6-7: Components and DDoS Workflow
    for s in slides_content[6:8]:
        add_bullet_slide(prs, s["title"], s["content"])
    
    # Add Predictive Scaling Model (slide 8)
    add_bullet_slide(prs, slides_content[8]["title"], slides_content[8]["content"])
    
    # Add equation slide if LaTeX tools available
    eq = r"\hat{y}_{t+1:t+12} = f(\mathbf{x}_{t-59:t})"
    if latex_to_png(eq, "equation.png"):
        add_image_slide(prs, "Predictive Scaling Equation", "equation.png")
    
    # Add slides 9-10: Key Features and Experimental Setup
    for s in slides_content[9:11]:
        add_bullet_slide(prs, s["title"], s["content"])
    
    # Add performance comparison table
    table_data = [
        ["Scaler", "Avg Pods", "Timeout Rate", "Scale Lag"],
        ["KEDA (Prometheus)", "7.2", "18.7%", "42s"],
        ["PredictKube", "8.1", "12.6%", "19s"],
        ["Improvement", "+12.5%", "-32.6%", "-54.8%"]
    ]
    add_table_slide(prs, "Scaling Performance at 500 RPS", table_data, [2, 2, 2, 2])
    
    # Add remaining slides (11-16)
    for s in slides_content[11:]:
        add_bullet_slide(prs, s["title"], s["content"])
    
    # Save final PPTX
    output_file = "Cloud_Center_Security_Complete.pptx"
    prs.save(output_file)
    
    # Clean up temporary image files
    for img in ["diagram.png", "equation.png"]:
        if os.path.exists(img):
            os.remove(img)
    
    print(f"✅ Generated {output_file}")
    print(f"   - {len(slides_content)} content slides from ppt.py")
    print("   - Enhanced with table, equation, and diagram from full_ppt.py")
    print(f"   - Total slides: ~{len(slides_content) + 3} (including visual enhancements)")

if __name__ == "__main__":
    main()
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# ---------------------------
# CONFIG: Text slides content
# ---------------------------
slides_content = [
    {
        "title": "Cloud Center Security: Enhancing Kubernetes Security",
        "subtitle": "Cilium, eBPF, and DDoS Mitigation\nZhijun Jiang, Amin Milani Fard, NYIT, Vancouver, Canada"
    },
    {
        "title": "Motivation",
        "content": [
            "Kubernetes widely used for container orchestration.",
            "Cloud services are increasingly targeted by DDoS attacks.",
            "Current autoscaling is reactive and lacks traffic differentiation.",
            "Need security-aware, predictive autoscaling."
        ]
    },
    {
        "title": "Problem Statement",
        "content": [
            "Reactive scaling triggers on all traffic → may scale up for attacks.",
            "Network and security policies lag behind scaling events.",
            "Observability is fragmented → difficult to detect malicious patterns."
        ]
    },
    {
        "title": "Objectives",
        "content": [
            "Enhance Kubernetes autoscaling with security awareness.",
            "Integrate Cilium + eBPF for fast, programmable networking.",
            "Implement DDoS mitigation pipeline.",
            "Enable predictive scaling using ML traffic forecasts.",
            "Provide observability via Hubble + Grafana dashboards."
        ]
    },
    {
        "title": "Architecture Overview",
        "content": [
            "Edge Load Balancer / Cloud Scrubber",
            "XDP eBPF pre-filter → drop/rate-limit malicious traffic",
            "Cilium dataplane → enforces L3/L4/L7 policies",
            "Hubble → observability and flow logs",
            "Security-aware autoscaler → ML + Prometheus + policies"
        ]
    },
    {
        "title": "Components",
        "content": [
            "XDP / eBPF: Early packet filtering, drops malicious traffic",
            "Cilium: Pod-to-pod networking, policy enforcement",
            "Hubble: Observability, flow logs, metrics",
            "Predictive Autoscaler: Forecast traffic, scale pods securely",
            "DDoS Engine: Attack detection and mitigation actions"
        ]
    },
    {
        "title": "DDoS Mitigation Workflow",
        "content": [
            "1. Traffic enters cluster",
            "2. XDP pre-filter drops attack traffic",
            "3. Cilium enforces network policies",
            "4. Hubble observes flows → alerts DDoS engine",
            "5. Security-aware autoscaler scales only legitimate traffic"
        ]
    },
    {
        "title": "Predictive Scaling Model",
        "content": [
            "LSTM-based time series prediction",
            "Input: last 60 time steps of traffic metrics",
            "Output: next 12-step traffic forecast",
            "Integrates with Prometheus metrics and autoscaler",
            "Equation: ŷ_{t+1:t+12} = f(x_{t-59:t})"
        ]
    },
    {
        "title": "Key Features",
        "content": [
            "Real-time security-aware scaling",
            "Fast eBPF-based load balancing via Cilium",
            "Observability: per-pod traffic flow, status code tracking",
            "Predictive: anticipates spikes to avoid overloading",
            "Integrated DDoS mitigation"
        ]
    },
    {
        "title": "Experimental Setup",
        "content": [
            "Cluster: AWS EKS, 4 nodes (m5.xlarge)",
            "Workload: Go microservice, 5-min startup",
            "Traffic: Vegeta simulated (legitimate + attack mix)",
            "Monitoring: Prometheus + Grafana + Hubble",
            "Autoscalers: KEDA + PredictKube + custom ML controller"
        ]
    },
    {
        "title": "Results: Connection Stability",
        "content": [
            "Connection drops reduced by 72% with Cilium + autoscaler",
            "Predictive autoscaler reduced scaling lag by 54.8%",
            "ML forecast prevented unnecessary pod scaling during DDoS"
        ]
    },
    {
        "title": "Results: DDoS Mitigation",
        "content": [
            "XDP pre-filter dropped malicious requests immediately",
            "Cilium network policies restricted attack propagation",
            "Observability via Hubble → real-time flow detection"
        ]
    },
    {
        "title": "Related Work",
        "content": [
            "Kubernetes autoscaling: HPA, KEDA",
            "eBPF and Cilium for cloud security",
            "DDoS detection and mitigation in cloud",
            "Predictive scaling for container workloads"
        ]
    },
    {
        "title": "Lessons Learned",
        "content": [
            "Observability is critical → Grafana + Hubble",
            "Predictive scaling improves resource efficiency",
            "DDoS mitigation must integrate with scaling → avoid waste",
            "Cilium + eBPF ensures instant networking convergence"
        ]
    },
    {
        "title": "Future Work",
        "content": [
            "Hybrid predictive + reactive security-aware controller",
            "Multi-cloud deployment for resilience",
            "Chaos engineering to test autoscaler under attack",
            "Integrate anomaly detection for novel attacks"
        ]
    },
    {
        "title": "Conclusion",
        "content": [
            "Cloud center security requires integrated network + scaling + observability",
            "Cilium + eBPF + ML autoscaler improves performance and mitigates DDoS",
            "Security-aware scaling ensures efficient resource use",
            "Architecture is extendable and observable"
        ]
    },
    {
        "title": "References",
        "content": [
            "Kubernetes Event-Driven Autoscaling (KEDA): https://keda.sh/",
            "Cilium: eBPF Networking for Kubernetes: https://cilium.io/",
            "Hubble: Observability for Cilium: https://cilium.io/hubble/",
            "DDoS Attack Detection & Mitigation Survey, IEEE Access, 2022",
            "PredictKube: ML for Kubernetes Autoscaling: https://predictkube.com/",
            "eBPF Literature: “BPF Performance and Applications”, ACM 2020"
        ]
    }
]

# ---------------------------
# CREATE PRESENTATION
# ---------------------------
prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = slides_content[0]["title"]
slide.placeholders[1].text = slides_content[0]["subtitle"]

# Add other slides
for s in slides_content[1:]:
    slide_layout = prs.slide_layouts[1]  # Title + content layout
    slide_obj = prs.slides.add_slide(slide_layout)
    slide_obj.shapes.title.text = s["title"]
    body = slide_obj.placeholders[1].text_frame
    if "content" in s:
        # Clear default paragraph
        body.clear()
        for line in s["content"]:
            p = body.add_paragraph()
            p.text = line
            p.level = 0  # bullet level

# Save PPTX
prs.save("Cloud_Center_Security_Presentation.pptx")
print("Presentation generated: Cloud_Center_Security_Presentation.pptx")
